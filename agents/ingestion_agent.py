# agents/ingestion_agent.py
import fitz  # PyMuPDF
import pandas as pd
import docx
import pptx
import os

class IngestionAgent:
    def __init__(self):
        pass

    def ingest(self, uploaded_files):
        all_chunks = []
        for file in uploaded_files:
            file_type = file.name.split('.')[-1].lower()
            if file_type == 'pdf':
                chunks = self._parse_pdf(file)
            elif file_type == 'docx':
                chunks = self._parse_docx(file)
            elif file_type == 'pptx':
                chunks = self._parse_pptx(file)
            elif file_type == 'csv':
                chunks = self._parse_csv(file)
            elif file_type in ['txt', 'md']:
                chunks = self._parse_txt(file)
            else:
                chunks = ["Unsupported file format"]
            all_chunks.extend(chunks)
        return all_chunks

    def _parse_pdf(self, file):
        doc = fitz.open(stream=file.read(), filetype="pdf")
        return [page.get_text() for page in doc]

    def _parse_docx(self, file):
        doc = docx.Document(file)
        return [para.text for para in doc.paragraphs if para.text.strip()]

    def _parse_pptx(self, file):
        prs = pptx.Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return text_runs

    def _parse_csv(self, file):
        df = pd.read_csv(file)
        return df.astype(str).apply(lambda row: ', '.join(row), axis=1).tolist()

    def _parse_txt(self, file):
        return file.read().decode("utf-8").splitlines()
